#!/usr/bin/env python3
"""
Build a landscape report from the output of export_comet_runs.py.

Usage:
    python diffusion_template/tools/comet/build_comet_report_pdf.py \
        --config diffusion_template/tools/comet/comet_pdf_config_template.json \
        --output diffusion_template/comet_data/comet_report.pdf

If the config omits "runs" or sets it to an empty list, all runs from the
export JSON are included in the report in export order.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import math
import re
import tempfile
import textwrap
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.gridspec import GridSpec
from PIL import Image, ImageOps

SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATE_ROOT = SCRIPT_DIR.parents[1]
DEFAULT_CONFIG_PATH = SCRIPT_DIR / "comet_pdf_config_template.json"
DEFAULT_EXPORT_JSON = TEMPLATE_ROOT / "comet_data" / "comet_runs_export.json"
LANDSCAPE_A4 = (11.69, 8.27)
IMG_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
PLACEHOLDER_RUN_ID = "REPLACE_WITH_COMET_EXPERIMENT_KEY"
MAX_CHARTS_PER_PAGE = 6
MARKDOWN_PAGE_RE = re.compile(r"^<!--\s*report-page:\s*(.+?)\s*-->$")
MARKDOWN_LAYOUT_RE = re.compile(r"^<!--\s*layout:\s*(.+?)\s*-->$")
MARKDOWN_COLUMN_RE = re.compile(r"^<!--\s*column:\s*(.+?)\s*-->$")
MARKDOWN_IMAGE_RE = re.compile(r"^!\[([^]]+)\]\(([^)]+)\)\s*$")
MARKDOWN_NUMBERED_RE = re.compile(r"^(\d+)\.\s+(.+)$")
PREFERRED_METRICS = [
    "train_loss",
    "loss",
    "val_loss",
    "train/loss",
    "val/loss",
    "accuracy",
    "val_accuracy",
    "train_accuracy",
    "general/steps_per_sec",
]
DEFAULT_KEY_HYPERPARAMETERS = [
    "step_shown",
    "model.weight_dtype",
    "model.rank",
    "dataloaders.train.batch_size",
    "branched_attn_weight_mode",
    "branched_attn_new_weight_kind",
    "train_branched_ca_lora",
    "lr_for_lora",
    "loss_kind",
    "lambda_face",
    "strict_face_routing",
    "trainer.masked_loss_step",
]


class ConfigError(ValueError):
    """Raised when the PDF configuration is invalid."""


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create a landscape report comparing exported Comet images and metrics."
    )
    parser.add_argument(
        "--config",
        type=Path,
        required=True,
        help=f"Path to the JSON config file. Template: {DEFAULT_CONFIG_PATH}",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="Path to the report file that will be created. Supports .pdf and .pptx.",
    )
    parser.add_argument(
        "--output-pdf",
        type=Path,
        default=None,
        help="Deprecated alias for --output.",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=200,
        help="Render DPI for report pages. Default: 200",
    )
    parser.add_argument(
        "--image-max-side",
        type=int,
        default=768,
        help="Maximum side length used when loading images. Default: 768",
    )
    parser.add_argument(
        "--image-dpi-percent",
        type=float,
        default=None,
        help=(
            "Optional override for the JSON image_dpi_percent setting. "
            "Examples: 100 preserves original size, 50 uses half size."
        ),
    )
    args = parser.parse_args()
    if args.output is None and args.output_pdf is None:
        parser.error("one of --output or --output-pdf is required")
    if args.output is not None and args.output_pdf is not None:
        parser.error("use only one of --output or --output-pdf")
    return args


def load_json(path: Path) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError as exc:
        raise ConfigError(f"JSON file does not exist: {path}") from exc
    except json.JSONDecodeError as exc:
        raise ConfigError(f"Invalid JSON file: {path}") from exc


def resolve_path(base_dir: Path, raw_path: str | None, fallback: Path | None = None) -> Path:
    if raw_path:
        path = Path(raw_path).expanduser()
        if not path.is_absolute():
            path = (base_dir / path).resolve()
        return path
    if fallback is None:
        raise ConfigError("A required path value is missing.")
    return fallback.resolve()


def parse_positive_int(value: Any, field_name: str) -> int:
    if isinstance(value, bool):
        raise ConfigError(f"{field_name} must be a positive integer")
    try:
        parsed = int(value)
    except (TypeError, ValueError) as exc:
        raise ConfigError(f"{field_name} must be a positive integer") from exc
    if parsed <= 0:
        raise ConfigError(f"{field_name} must be a positive integer")
    return parsed


def parse_bool(value: Any, field_name: str) -> bool:
    if isinstance(value, bool):
        return value
    raise ConfigError(f"{field_name} must be true or false")


def parse_optional_positive_number(value: Any, field_name: str) -> float | None:
    if value is None:
        return None
    if isinstance(value, bool):
        raise ConfigError(f"{field_name} must be a positive number")
    try:
        parsed = float(value)
    except (TypeError, ValueError) as exc:
        raise ConfigError(f"{field_name} must be a positive number") from exc
    if parsed <= 0:
        raise ConfigError(f"{field_name} must be a positive number")
    return parsed


def parse_nonnegative_number(value: Any, field_name: str) -> float:
    if isinstance(value, bool):
        raise ConfigError(f"{field_name} must be a non-negative number")
    try:
        parsed = float(value)
    except (TypeError, ValueError) as exc:
        raise ConfigError(f"{field_name} must be a non-negative number") from exc
    if not math.isfinite(parsed) or parsed < 0:
        raise ConfigError(f"{field_name} must be a non-negative number")
    return parsed


def parse_output_format(value: Any) -> str:
    if value is None:
        return "pdf"
    text = str(value).strip().lower()
    if text in {"pdf", "ppt", "pptx"}:
        return "pptx" if text in {"ppt", "pptx"} else "pdf"
    raise ConfigError("output_format must be 'pdf' or 'pptx'")


def natural_sort_key(text: str) -> list[Any]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", text)]


def chunked(items: list[Any], size: int) -> list[list[Any]]:
    return [items[index : index + size] for index in range(0, len(items), size)]


def wrap_label(text: str, width: int) -> str:
    normalized = " ".join(str(text).split())
    return "\n".join(textwrap.wrap(normalized, width=width)[:4]) if normalized else ""


def load_image_for_report(
    path: Path,
    max_side: int,
    image_dpi_percent: float | None,
    crop_bbox: list[float] | None = None,
    crop_padding_ratio: float = 0.0,
    bbox_coordinate_size: tuple[int, int] | None = None,
) -> Image.Image:
    image = Image.open(path).convert("RGB")
    if crop_bbox is not None:
        image = square_face_crop(
            image,
            crop_bbox,
            padding_ratio=crop_padding_ratio,
            bbox_coordinate_size=bbox_coordinate_size,
        )
    if image_dpi_percent is not None:
        scale = image_dpi_percent / 100.0
        new_size = (
            max(1, round(image.size[0] * scale)),
            max(1, round(image.size[1] * scale)),
        )
        if new_size != image.size:
            image = image.resize(new_size, Image.LANCZOS)
    elif max(image.size) > max_side:
        image.thumbnail((max_side, max_side), Image.LANCZOS)
    return image


def square_face_crop(
    image: Image.Image,
    bbox: list[float],
    padding_ratio: float,
    bbox_coordinate_size: tuple[int, int] | None,
) -> Image.Image:
    """Crop a square around a fixed owned-face bbox without re-detecting faces."""
    x0, y0, x1, y1 = bbox
    if bbox_coordinate_size is not None:
        coordinate_width, coordinate_height = bbox_coordinate_size
        x_scale = image.width / coordinate_width
        y_scale = image.height / coordinate_height
        x0, x1 = x0 * x_scale, x1 * x_scale
        y0, y1 = y0 * y_scale, y1 * y_scale

    if x0 < 0 or y0 < 0 or x1 > image.width or y1 > image.height:
        raise ConfigError(
            f"Face bbox {bbox} resolves outside image canvas {image.size}"
        )

    center_x = (x0 + x1) / 2.0
    center_y = (y0 + y1) / 2.0
    face_side = max(x1 - x0, y1 - y0)
    side = max(2, int(round(face_side * (1.0 + 2.0 * padding_ratio))))
    side = min(side, image.width, image.height)
    crop_x0 = int(round(center_x - side / 2.0))
    crop_y0 = int(round(center_y - side / 2.0))
    crop_x0 = min(max(crop_x0, 0), image.width - side)
    crop_y0 = min(max(crop_y0, 0), image.height - side)
    return image.crop((crop_x0, crop_y0, crop_x0 + side, crop_y0 + side))


def placeholder_image(size: tuple[int, int] = (512, 512)) -> Image.Image:
    image = Image.new("RGB", size, "#f2f2f2")
    return ImageOps.expand(image, border=2, fill="#cccccc")


def display_label_for_image(file_name: str) -> str:
    path = Path(file_name)
    label = path.stem if path.stem else path.name
    label = re.sub(r"\s*\(\d+\)$", "", label).strip()
    label = re.sub(r"_{1,}\d+$", "", label).strip()
    return label or path.stem or path.name


def canonical_image_key(file_name: str) -> str:
    # Comet image assets replace spaces with underscores, while per-image CSV
    # output_key values retain spaces. Normalize both sides before joining.
    return re.sub(r"\s", "_", display_label_for_image(file_name))


def is_mask_image_name(file_name: str) -> bool:
    raw_stem = Path(file_name).stem
    normalized_stem = display_label_for_image(file_name)
    return raw_stem.endswith("_mask") or normalized_stem.endswith("_mask")


def run_display_name(run_cfg: dict[str, Any], export_run: dict[str, Any]) -> str:
    return str(run_cfg.get("run_name") or export_run.get("name") or export_run.get("id"))


def normalize_string_list(value: Any, field_name: str) -> list[str]:
    if value is None:
        return []
    if not isinstance(value, list):
        raise ConfigError(f"'{field_name}' must be a list when provided.")
    return [str(item).strip() for item in value if str(item).strip()]


def normalize_object(value: Any, field_name: str) -> dict[str, Any]:
    if value is None:
        return {}
    if not isinstance(value, dict):
        raise ConfigError(f"'{field_name}' must be an object when provided.")
    return dict(value)


def normalize_face_closeups(
    value: Any,
    config_dir: Path,
) -> dict[str, Any]:
    raw = normalize_object(value, "face_closeups")
    enabled = parse_bool(raw.get("enabled", False), "face_closeups.enabled")
    normalized: dict[str, Any] = {
        "enabled": enabled,
        "padding_ratio": parse_nonnegative_number(
            raw.get("padding_ratio", 0.45),
            "face_closeups.padding_ratio",
        ),
        "bbox_field": str(raw.get("bbox_field") or "face_crop_new").strip(),
        "title": str(raw.get("title") or "Comet Run Face Region Comparison").strip(),
        "require_exact_keys": parse_bool(
            raw.get("require_exact_keys", True),
            "face_closeups.require_exact_keys",
        ),
        "bbox_json": None,
        "expected_sha256": None,
        "bbox_coordinate_size": None,
    }
    if not normalized["bbox_field"]:
        raise ConfigError("face_closeups.bbox_field must not be empty")
    if not normalized["title"]:
        raise ConfigError("face_closeups.title must not be empty")

    coordinate_size = raw.get("bbox_coordinate_size")
    if coordinate_size is not None:
        if not isinstance(coordinate_size, list) or len(coordinate_size) != 2:
            raise ConfigError(
                "face_closeups.bbox_coordinate_size must be [width, height]"
            )
        normalized["bbox_coordinate_size"] = (
            parse_positive_int(
                coordinate_size[0], "face_closeups.bbox_coordinate_size[0]"
            ),
            parse_positive_int(
                coordinate_size[1], "face_closeups.bbox_coordinate_size[1]"
            ),
        )

    expected_sha256 = str(raw.get("expected_sha256") or "").strip().lower()
    if expected_sha256:
        if not re.fullmatch(r"[0-9a-f]{64}", expected_sha256):
            raise ConfigError(
                "face_closeups.expected_sha256 must be a 64-character hex digest"
            )
        normalized["expected_sha256"] = expected_sha256

    bbox_json_raw = raw.get("bbox_json")
    if enabled:
        if bbox_json_raw in (None, ""):
            raise ConfigError(
                "face_closeups.bbox_json is required when face close-ups are enabled"
            )
        bbox_json = resolve_path(config_dir, str(bbox_json_raw))
        if not bbox_json.is_file():
            raise ConfigError(f"Face bbox JSON does not exist: {bbox_json}")
        normalized["bbox_json"] = bbox_json
    return normalized


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_face_bbox_map(
    face_config: dict[str, Any],
    image_keys: list[str],
) -> dict[str, list[float]]:
    # 21 Aug 2026 - AICODE-NOTE: Face comparison pages use one sealed fixed-box
    # protocol for every run. Per-run face detection would hide ownership drift.
    if not face_config.get("enabled", False):
        return {}

    bbox_json = Path(face_config["bbox_json"])
    expected_sha256 = face_config.get("expected_sha256")
    actual_sha256 = file_sha256(bbox_json)
    if expected_sha256 and actual_sha256 != expected_sha256:
        raise ConfigError(
            f"Face bbox JSON SHA-256 mismatch: expected {expected_sha256}, "
            f"found {actual_sha256} at {bbox_json}"
        )

    payload = load_json(bbox_json)
    if not isinstance(payload, dict):
        raise ConfigError(f"Face bbox JSON must contain an object: {bbox_json}")
    bbox_field = str(face_config["bbox_field"])
    bbox_map: dict[str, list[float]] = {}
    for raw_key, raw_record in payload.items():
        if not isinstance(raw_record, dict):
            raise ConfigError(f"Invalid bbox record for {raw_key!r} in {bbox_json}")
        raw_bbox = raw_record.get(bbox_field)
        if not isinstance(raw_bbox, (list, tuple)) or len(raw_bbox) != 4:
            raise ConfigError(
                f"BBox record {raw_key!r} has no valid {bbox_field!r} field"
            )
        if any(isinstance(item, bool) for item in raw_bbox):
            raise ConfigError(f"Non-numeric face bbox for {raw_key!r}")
        try:
            bbox = [float(item) for item in raw_bbox]
        except (TypeError, ValueError) as exc:
            raise ConfigError(f"Non-numeric face bbox for {raw_key!r}") from exc
        x0, y0, x1, y1 = bbox
        if not all(math.isfinite(item) for item in bbox) or x1 <= x0 or y1 <= y0:
            raise ConfigError(f"Invalid face bbox for {raw_key!r}: {raw_bbox}")
        key = canonical_image_key(str(raw_key))
        if key in bbox_map:
            raise ConfigError(f"Duplicate canonical face bbox key {key!r}")
        bbox_map[key] = bbox

    image_key_set = set(image_keys)
    bbox_key_set = set(bbox_map)
    missing = sorted(image_key_set - bbox_key_set, key=natural_sort_key)
    extra = sorted(bbox_key_set - image_key_set, key=natural_sort_key)
    if missing or (face_config.get("require_exact_keys", True) and extra):
        raise ConfigError(
            "Image/bbox join mismatch "
            f"images={len(image_keys)} boxes={len(bbox_map)} "
            f"missing={missing[:5]} extra={extra[:5]}"
        )
    return bbox_map


def select_numeric_metrics(
    selected_runs: list[dict[str, Any]],
    configured_metrics: list[str] | None,
) -> list[str]:
    if configured_metrics:
        return [metric for metric in configured_metrics if metric]

    run_metric_sets = []
    for run in selected_runs:
        numeric_metrics: set[str] = set()
        for metric_name, entries in run["export_run"].get("metrics", {}).items():
            if extract_metric_series(entries):
                numeric_metrics.add(metric_name)
        run_metric_sets.append(numeric_metrics)

    if not run_metric_sets:
        return []

    common_metrics = set.intersection(*run_metric_sets) if run_metric_sets else set()
    candidate_metrics = common_metrics or set.union(*run_metric_sets)

    def metric_priority(metric_name: str) -> tuple[int, str]:
        if metric_name in PREFERRED_METRICS:
            return (PREFERRED_METRICS.index(metric_name), metric_name.lower())
        return (len(PREFERRED_METRICS) + 1, metric_name.lower())

    return sorted(candidate_metrics, key=metric_priority)[:MAX_CHARTS_PER_PAGE]


def extract_metric_series(entries: Any) -> list[tuple[float, float]]:
    series: list[tuple[float, float]] = []
    if not isinstance(entries, list):
        return series
    for entry in entries:
        if not isinstance(entry, dict):
            continue
        step = entry.get("step")
        value = entry.get("value")
        if isinstance(step, bool) or isinstance(value, bool):
            continue
        try:
            step_num = float(step)
            value_num = float(value)
        except (TypeError, ValueError):
            continue
        series.append((step_num, value_num))
    series.sort(key=lambda pair: pair[0])
    return series


def load_pdf_config(path: Path) -> tuple[dict[str, Any], Path]:
    raw_config = load_json(path)
    if not isinstance(raw_config, dict):
        raise ConfigError("PDF config must be a JSON object.")

    config_dir = path.resolve().parent
    runs = raw_config.get("runs", [])
    if runs is None:
        runs = []
    if not isinstance(runs, list):
        raise ConfigError("'runs' must be a list when provided.")

    normalized_runs: list[dict[str, Any]] = []
    for index, run in enumerate(runs, start=1):
        if not isinstance(run, dict):
            raise ConfigError(f"runs[{index}] must be an object.")
        run_id = str(run.get("run_id", "")).strip()
        if not run_id:
            raise ConfigError(f"runs[{index}].run_id is required.")
        if run_id == PLACEHOLDER_RUN_ID:
            raise ConfigError(
                f"runs[{index}].run_id still uses the template placeholder {PLACEHOLDER_RUN_ID}."
            )
        run_name = run.get("run_name")
        hyperparameter_overrides = normalize_object(
            run.get("hyperparameter_overrides"),
            f"runs[{index}].hyperparameter_overrides",
        )
        normalized_runs.append(
            {
                "run_id": run_id,
                "run_name": None if run_name in (None, "") else str(run_name).strip(),
                "hyperparameter_overrides": hyperparameter_overrides,
            }
        )

    configured_metrics = normalize_string_list(raw_config.get("key_metrics"), "key_metrics")
    if not configured_metrics:
        configured_metrics = None
    configured_hyperparameters = normalize_string_list(
        raw_config.get("key_hyperparameters", DEFAULT_KEY_HYPERPARAMETERS),
        "key_hyperparameters",
    )
    hyperparameter_labels = {
        str(key): str(value)
        for key, value in normalize_object(
            raw_config.get("hyperparameter_labels"), "hyperparameter_labels"
        ).items()
    }
    per_image_metric = normalize_object(
        raw_config.get("per_image_metric"), "per_image_metric"
    )
    group_average_tables = normalize_object(
        raw_config.get("group_average_tables"), "group_average_tables"
    )
    metric_point_labels = normalize_object(
        raw_config.get("metric_point_labels"), "metric_point_labels"
    )
    face_closeups = normalize_face_closeups(
        raw_config.get("face_closeups"),
        config_dir,
    )
    markdown_source_raw = raw_config.get("markdown_source")
    markdown_source = None
    if markdown_source_raw not in (None, ""):
        markdown_source = resolve_path(config_dir, str(markdown_source_raw))
        if not markdown_source.is_file():
            raise ConfigError(f"Markdown source does not exist: {markdown_source}")

    normalized_config = {
        "export_json": resolve_path(config_dir, raw_config.get("export_json"), DEFAULT_EXPORT_JSON),
        "output_format": parse_output_format(raw_config.get("output_format", "pdf")),
        "max_columns": parse_positive_int(raw_config.get("max_columns", 4), "max_columns"),
        "max_rows": parse_positive_int(raw_config.get("max_rows", 3), "max_rows"),
        "ignore_mask": parse_bool(raw_config.get("ignore_mask", True), "ignore_mask"),
        "image_dpi_percent": parse_optional_positive_number(
            raw_config.get("image_dpi_percent"),
            "image_dpi_percent",
        ),
        "run_name_max_chars_per_line": parse_positive_int(
            raw_config.get("run_name_max_chars_per_line", 16),
            "run_name_max_chars_per_line",
        ),
        "key_metrics": configured_metrics,
        "key_hyperparameters": configured_hyperparameters,
        "hyperparameter_labels": hyperparameter_labels,
        "per_image_metric": per_image_metric,
        "group_average_tables": group_average_tables,
        "metric_point_labels": metric_point_labels,
        "face_closeups": face_closeups,
        "markdown_source": markdown_source,
        "runs": normalized_runs,
    }
    return normalized_config, config_dir


def prepare_selected_runs(
    config_runs: list[dict[str, Any]],
    export_payload: dict[str, Any],
    export_json_path: Path,
    ignore_mask: bool,
) -> tuple[list[dict[str, Any]], Path]:
    export_runs = export_payload.get("runs")
    if not isinstance(export_runs, list):
        raise ConfigError(f"'runs' is missing or invalid in export JSON: {export_json_path}")

    ordered_export_runs: list[dict[str, Any]] = []
    run_by_id: dict[str, dict[str, Any]] = {}
    for run in export_runs:
        if isinstance(run, dict) and run.get("id"):
            run_id = str(run["id"])
            ordered_export_runs.append(run)
            run_by_id[run_id] = run

    images_root = export_payload.get("output_dir")
    if not images_root:
        raise ConfigError(f"'output_dir' is missing in export JSON: {export_json_path}")
    image_root_dir = Path(images_root).expanduser().resolve()

    selected_runs: list[dict[str, Any]] = []
    if not config_runs:
        for export_run in ordered_export_runs:
            selected_runs.append(
                {
                    "id": export_run["id"],
                    "name": str(export_run.get("name") or export_run.get("id")),
                    "export_run": export_run,
                    "image_map": build_run_image_map(export_run, image_root_dir, ignore_mask),
                    "hyperparameter_overrides": {},
                }
            )
        if not selected_runs:
            raise ConfigError(f"No runs were found in export JSON: {export_json_path}")
        return selected_runs, image_root_dir

    missing_ids: list[str] = []
    for run_cfg in config_runs:
        export_run = run_by_id.get(run_cfg["run_id"])
        if export_run is None:
            missing_ids.append(run_cfg["run_id"])
            continue
        selected_runs.append(
            {
                "id": export_run["id"],
                "name": run_display_name(run_cfg, export_run),
                "export_run": export_run,
                "image_map": build_run_image_map(export_run, image_root_dir, ignore_mask),
                "hyperparameter_overrides": run_cfg.get("hyperparameter_overrides", {}),
            }
        )

    if missing_ids:
        raise ConfigError(f"Run IDs not found in export JSON: {', '.join(missing_ids)}")
    if not selected_runs:
        raise ConfigError("No runs were selected for the PDF report.")

    return selected_runs, image_root_dir


def build_run_image_map(
    export_run: dict[str, Any],
    image_root_dir: Path,
    ignore_mask: bool,
) -> dict[str, dict[str, Any]]:
    image_map: dict[str, dict[str, Any]] = {}
    for image_info in export_run.get("downloaded_images", []):
        if not isinstance(image_info, dict):
            continue
        file_name = str(image_info.get("file_name") or "").strip()
        if not file_name:
            continue
        if ignore_mask and is_mask_image_name(file_name):
            continue
        key = canonical_image_key(file_name)
        relative_path = image_info.get("saved_path")
        if relative_path:
            image_path = (image_root_dir / relative_path).resolve()
        else:
            output_folder = export_run.get("output_folder", "")
            image_path = (image_root_dir / str(output_folder) / file_name).resolve()
        image_map[key] = {
            "file_name": file_name,
            "display_name": display_label_for_image(file_name),
            "path": image_path,
        }
    return image_map


def attach_per_image_metric_data(
    selected_runs: list[dict[str, Any]],
    image_root_dir: Path,
    metric_config: dict[str, Any],
) -> None:
    if not metric_config.get("enabled", False):
        return

    metric_column = str(metric_config.get("column") or "id_sim")
    for run in selected_runs:
        tables = run["export_run"].get("downloaded_tables", [])
        candidates = [
            table
            for table in tables
            if isinstance(table, dict) and table.get("kind") == "per_image_id"
        ]
        if len(candidates) != 1:
            raise ConfigError(
                f"{run['name']}: expected exactly one per-image ID table, "
                f"found {len(candidates)}"
            )
        saved_path = candidates[0].get("saved_path")
        table_path = (image_root_dir / str(saved_path)).resolve()
        if not table_path.is_file():
            raise ConfigError(f"{run['name']}: per-image table is missing: {table_path}")

        with table_path.open(newline="", encoding="utf-8") as handle:
            rows = list(csv.DictReader(handle))
        if not rows or "output_key" not in rows[0] or metric_column not in rows[0]:
            raise ConfigError(
                f"{run['name']}: table must contain output_key and {metric_column!r}"
            )

        metric_map: dict[str, float] = {}
        for row in rows:
            key = canonical_image_key(str(row.get("output_key") or ""))
            try:
                value = float(row[metric_column])
            except (TypeError, ValueError, KeyError):
                raise ConfigError(
                    f"{run['name']}: non-numeric {metric_column!r} for {key!r}"
                )
            if key in metric_map:
                raise ConfigError(f"{run['name']}: duplicate per-image key {key!r}")
            metric_map[key] = value

        missing = sorted(set(run["image_map"]) - set(metric_map), key=natural_sort_key)
        extra = sorted(set(metric_map) - set(run["image_map"]), key=natural_sort_key)
        if missing or extra or len(rows) != len(run["image_map"]):
            raise ConfigError(
                f"{run['name']}: image/metric join mismatch "
                f"images={len(run['image_map'])} rows={len(rows)} "
                f"missing={missing[:5]} extra={extra[:5]}"
            )
        run["per_image_metric_map"] = metric_map
        run["per_image_metric_rows"] = rows


def prompt_group_label(row: dict[str, str]) -> str:
    output_key = display_label_for_image(str(row.get("output_key") or ""))
    identity = str(row.get("identity") or "").strip()
    suffix = f"_{identity}" if identity else ""
    if suffix and output_key.lower().endswith(suffix.lower()):
        output_key = output_key[: -len(suffix)]
    return " ".join(output_key.replace("_", " ").split()) or "Unknown"


def parse_markdown_pages(path: Path) -> list[dict[str, Any]]:
    pages: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    current_column = "main"
    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.rstrip()
        page_match = MARKDOWN_PAGE_RE.match(line)
        if page_match:
            current = {
                "name": page_match.group(1).strip(),
                "layout": "text",
                "columns": {"main": []},
                "source_dir": path.parent,
            }
            pages.append(current)
            current_column = "main"
            continue
        if current is None:
            if line.strip():
                raise ConfigError(
                    f"Markdown content before first report-page marker in {path}"
                )
            continue
        layout_match = MARKDOWN_LAYOUT_RE.match(line)
        if layout_match:
            current["layout"] = layout_match.group(1).strip()
            continue
        column_match = MARKDOWN_COLUMN_RE.match(line)
        if column_match:
            current_column = column_match.group(1).strip()
            current["columns"].setdefault(current_column, [])
            continue
        current["columns"].setdefault(current_column, []).append(line)
    return pages


def markdown_title(columns: dict[str, list[str]], fallback: str) -> str:
    for lines in columns.values():
        for line in lines:
            if line.startswith("# "):
                return line[2:].strip()
    return fallback


def render_markdown_column(
    ax: plt.Axes,
    lines: list[str],
    wrap_width: int,
) -> None:
    ax.axis("off")
    blocks: list[tuple[str, str]] = []
    paragraph: list[str] = []

    def flush_paragraph() -> None:
        if paragraph:
            blocks.append(("paragraph", " ".join(paragraph)))
            paragraph.clear()

    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith("# "):
            flush_paragraph()
            continue
        if stripped.startswith("## "):
            flush_paragraph()
            blocks.append(("heading", stripped[3:].strip()))
        elif stripped.startswith("- "):
            flush_paragraph()
            blocks.append(("bullet", stripped[2:].strip()))
        elif stripped.startswith("$$") and stripped.endswith("$$"):
            flush_paragraph()
            blocks.append(("formula", "$" + stripped[2:-2].strip() + "$"))
        else:
            paragraph.append(stripped)
    flush_paragraph()

    y = 0.985
    for kind, value in blocks:
        if kind == "heading":
            y -= 0.015
            ax.text(0.0, y, value, ha="left", va="top", fontsize=11.2, weight="bold")
            y -= 0.052
            continue
        if kind == "formula":
            ax.text(0.03, y, value, ha="left", va="top", fontsize=10.0)
            y -= 0.068
            continue
        prefix = "- " if kind == "bullet" else ""
        first_width = max(20, wrap_width - len(prefix))
        wrapped = textwrap.wrap(value, width=first_width) or [""]
        rendered = prefix + wrapped[0]
        if len(wrapped) > 1:
            rendered += "\n" + "\n".join(
                ("  " if kind == "bullet" else "") + item for item in wrapped[1:]
            )
        ax.text(
            0.0,
            y,
            rendered,
            ha="left",
            va="top",
            fontsize=8.4,
            linespacing=1.20,
            color="#202020",
        )
        y -= 0.034 * len(wrapped) + 0.016
    if y < -0.01:
        raise ConfigError("Markdown architecture page content exceeds its column")


def render_architecture_markdown_page(
    writer: "ReportWriter",
    page: dict[str, Any],
    dpi: int,
    page_number: int,
) -> int:
    columns = page["columns"]
    title = markdown_title(columns, page["name"])
    fig = plt.figure(figsize=LANDSCAPE_A4)
    fig.text(0.025, 0.975, title, ha="left", va="top", fontsize=14, weight="bold")
    left = fig.add_axes([0.03, 0.07, 0.54, 0.84])
    right = fig.add_axes([0.60, 0.07, 0.37, 0.84])
    render_markdown_column(left, columns.get("left", columns.get("main", [])), 83)
    render_markdown_column(right, columns.get("right", []), 55)
    add_page_number(fig, page_number)
    writer.save_figure(fig, dpi=dpi)
    plt.close(fig)
    return page_number + 1


def render_code_markdown_column(
    ax: plt.Axes,
    lines: list[str],
    wrap_width: int,
) -> None:
    ax.axis("off")
    blocks: list[tuple[str, str]] = []
    paragraph: list[str] = []
    code_lines: list[str] = []
    in_code = False

    def flush_paragraph() -> None:
        if paragraph:
            blocks.append(("paragraph", " ".join(paragraph)))
            paragraph.clear()

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```"):
            if in_code:
                blocks.append(("code", "\n".join(code_lines)))
                code_lines.clear()
                in_code = False
            else:
                flush_paragraph()
                in_code = True
            continue
        if in_code:
            code_lines.append(line.rstrip())
            continue
        if not stripped or stripped.startswith("# "):
            flush_paragraph()
        elif stripped.startswith("## "):
            flush_paragraph()
            blocks.append(("heading", stripped[3:].strip()))
        elif stripped.startswith("- "):
            flush_paragraph()
            blocks.append(("bullet", stripped[2:].strip()))
        else:
            paragraph.append(stripped)
    if in_code:
        raise ConfigError("Unclosed fenced code block in Markdown report page")
    flush_paragraph()

    y = 0.99
    for kind, value in blocks:
        if kind == "heading":
            y -= 0.01
            ax.text(0.0, y, value, ha="left", va="top", fontsize=10.6, weight="bold")
            y -= 0.047
        elif kind == "code":
            line_count = value.count("\n") + 1
            ax.text(
                0.012,
                y,
                value,
                ha="left",
                va="top",
                fontsize=6.45,
                family="monospace",
                linespacing=1.18,
                color="#102030",
                bbox={
                    "boxstyle": "round,pad=0.55",
                    "facecolor": "#f3f6f9",
                    "edgecolor": "#9aa7b5",
                    "linewidth": 0.65,
                },
            )
            y -= 0.0255 * line_count + 0.030
        else:
            prefix = "- " if kind == "bullet" else ""
            wrapped = textwrap.wrap(value, width=wrap_width) or [""]
            rendered = prefix + wrapped[0]
            if len(wrapped) > 1:
                rendered += "\n" + "\n".join(
                    ("  " if kind == "bullet" else "") + item
                    for item in wrapped[1:]
                )
            ax.text(
                0.0,
                y,
                rendered,
                ha="left",
                va="top",
                fontsize=8.0,
                linespacing=1.18,
            )
            y -= 0.032 * len(wrapped) + 0.014
    if y < -0.01:
        raise ConfigError("Markdown code page content exceeds its column")


def render_code_markdown_page(
    writer: "ReportWriter",
    page: dict[str, Any],
    dpi: int,
    page_number: int,
) -> int:
    columns = page["columns"]
    title = markdown_title(columns, page["name"])
    fig = plt.figure(figsize=LANDSCAPE_A4)
    fig.text(0.025, 0.975, title, ha="left", va="top", fontsize=14, weight="bold")
    left = fig.add_axes([0.025, 0.06, 0.46, 0.86])
    right = fig.add_axes([0.515, 0.06, 0.46, 0.86])
    render_code_markdown_column(left, columns.get("left", columns.get("main", [])), 70)
    render_code_markdown_column(right, columns.get("right", []), 70)
    add_page_number(fig, page_number)
    writer.save_figure(fig, dpi=dpi)
    plt.close(fig)
    return page_number + 1


def render_references_prompts_markdown_page(
    writer: "ReportWriter",
    page: dict[str, Any],
    dpi: int,
    page_number: int,
) -> int:
    columns = page["columns"]
    lines = [line for group in columns.values() for line in group]
    title = markdown_title(columns, page["name"])
    images: list[tuple[str, Path]] = []
    prompts: list[tuple[int, str]] = []
    for line in lines:
        image_match = MARKDOWN_IMAGE_RE.match(line.strip())
        if image_match:
            path = Path(image_match.group(2)).expanduser()
            if not path.is_absolute():
                path = (page["source_dir"] / path).resolve()
            if not path.is_file():
                raise ConfigError(f"Markdown reference image is missing: {path}")
            images.append((image_match.group(1).strip(), path))
            continue
        prompt_match = MARKDOWN_NUMBERED_RE.match(line.strip())
        if prompt_match:
            prompts.append((int(prompt_match.group(1)), prompt_match.group(2).strip()))
    if len(images) != 8 or len(prompts) != 12:
        raise ConfigError(
            f"Reference/prompt page requires exactly 8 images and 12 prompts; "
            f"found {len(images)} and {len(prompts)}"
        )

    fig = plt.figure(figsize=LANDSCAPE_A4)
    fig.text(0.025, 0.975, title, ha="left", va="top", fontsize=14, weight="bold")
    fig.text(0.025, 0.925, "Reference identities", ha="left", va="top", fontsize=11.5, weight="bold")
    fig.text(0.50, 0.925, "Prompt templates", ha="left", va="top", fontsize=11.5, weight="bold")

    cell_w, cell_h = 0.215, 0.182
    for index, (name, path) in enumerate(images):
        row, col = divmod(index, 2)
        ax = fig.add_axes([0.025 + col * 0.225, 0.73 - row * 0.205, cell_w, cell_h])
        image = Image.open(path).convert("RGB")
        ax.imshow(image)
        ax.set_xticks([])
        ax.set_yticks([])
        for spine in ax.spines.values():
            spine.set_color("#5b677a")
            spine.set_linewidth(0.8)
        ax.text(
            0.5,
            -0.055,
            name,
            transform=ax.transAxes,
            ha="center",
            va="top",
            fontsize=8.8,
            weight="bold",
        )

    prompt_ax = fig.add_axes([0.50, 0.085, 0.47, 0.80])
    prompt_ax.axis("off")
    y = 0.99
    for number, prompt in prompts:
        wrapped = textwrap.wrap(prompt, width=66) or [""]
        text_value = f"{number}. {wrapped[0]}"
        if len(wrapped) > 1:
            text_value += "\n   " + "\n   ".join(wrapped[1:])
        prompt_ax.text(
            0.0,
            y,
            text_value,
            ha="left",
            va="top",
            fontsize=8.35,
            linespacing=1.18,
        )
        y -= 0.069 + 0.028 * (len(wrapped) - 1)
    validation_contract = "\n".join(
        textwrap.wrap(
            "Validation contract: 8 references x 12 prompts x seed 0 = 96 images. "
            "<class> resolves to 'man img' or 'woman img'.",
            width=72,
        )
    )
    fig.text(
        0.50,
        0.048,
        validation_contract,
        ha="left",
        va="bottom",
        fontsize=7.8,
        color="#444444",
    )
    add_page_number(fig, page_number)
    writer.save_figure(fig, dpi=dpi)
    plt.close(fig)
    return page_number + 1


def render_markdown_pages(
    writer: "ReportWriter",
    markdown_source: Path | None,
    dpi: int,
    page_number: int,
) -> int:
    if markdown_source is None:
        return page_number
    for page in parse_markdown_pages(markdown_source):
        layout = str(page.get("layout") or "text")
        if layout == "architecture":
            page_number = render_architecture_markdown_page(
                writer, page, dpi, page_number
            )
        elif layout == "code":
            page_number = render_code_markdown_page(
                writer, page, dpi, page_number
            )
        elif layout == "references_prompts":
            page_number = render_references_prompts_markdown_page(
                writer, page, dpi, page_number
            )
        else:
            raise ConfigError(f"Unsupported Markdown report-page layout: {layout}")
    return page_number


def collect_image_keys(selected_runs: list[dict[str, Any]]) -> list[str]:
    image_keys: set[str] = set()
    for run in selected_runs:
        image_keys.update(run["image_map"].keys())
    return sorted(image_keys, key=natural_sort_key)


def add_page_number(fig: plt.Figure, page_number: int) -> None:
    fig.text(0.5, 0.015, f"Page {page_number}", ha="center", va="bottom", fontsize=9)


def step_display_value_for_run(run: dict[str, Any]) -> str:
    export_run = run["export_run"]
    selection = export_run.get("step_selection")
    requested = export_run.get("requested_step_number")
    resolved = export_run.get("resolved_step_number")

    if isinstance(selection, dict):
        requested = selection.get("requested_step_number", requested)
        resolved = selection.get("resolved_step_number", resolved)

    if resolved is None and requested is None:
        return "step unavailable"
    if resolved is None:
        return f"no image step for requested {requested}"
    if requested is None or requested == resolved:
        return f"step {resolved}"
    return f"step {resolved} (requested {requested})"


def step_note_for_run(run: dict[str, Any]) -> str:
    return f"{run['name']}: {step_display_value_for_run(run)}"


def add_step_footnote(fig: plt.Figure, runs: list[dict[str, Any]]) -> None:
    if not runs:
        return
    note = "Step notes: " + "; ".join(step_note_for_run(run) for run in runs)
    wrapped = "\n".join(textwrap.wrap(note, width=135))
    fig.text(0.04, 0.032, wrapped, ha="left", va="bottom", fontsize=8, color="#444444")


def render_missing_cell(ax: plt.Axes) -> None:
    ax.imshow(placeholder_image())
    ax.text(0.5, 0.5, "Missing", ha="center", va="center", fontsize=10, color="#555555")
    ax.set_xticks([])
    ax.set_yticks([])
    ax.set_frame_on(False)


def stringify_hyperparameter_value(value: Any) -> str:
    if value in (None, ""):
        return "n/a"
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        return f"{value:.6g}"
    if isinstance(value, (list, dict)):
        return json.dumps(value, ensure_ascii=True, sort_keys=True)
    text = str(value).strip()
    return text if text else "n/a"


def hyperparameter_value_for_run(run: dict[str, Any], hyperparameter_name: str) -> str:
    overrides = run.get("hyperparameter_overrides", {})
    if hyperparameter_name in overrides:
        return stringify_hyperparameter_value(overrides[hyperparameter_name])
    if hyperparameter_name == "step_shown":
        resolved = run["export_run"].get("resolved_step_number")
        return "step unavailable" if resolved is None else f"step {resolved}"

    export_run = run["export_run"]
    hyperparameters = export_run.get("hyperparameters", {})
    if not isinstance(hyperparameters, dict):
        return "n/a"

    return stringify_hyperparameter_value(hyperparameters.get(hyperparameter_name))


def render_hyperparameter_page(
    writer: "ReportWriter",
    selected_runs: list[dict[str, Any]],
    hyperparameter_names: list[str],
    hyperparameter_labels: dict[str, str],
    run_name_max_chars_per_line: int,
    dpi: int,
    page_number: int,
) -> int:
    if not hyperparameter_names:
        return page_number

    fig = plt.figure(figsize=LANDSCAPE_A4)
    fig.text(0.02, 0.982, "Key Hyperparameters", ha="left", va="top", fontsize=14, weight="bold")
    ax = fig.add_axes([0.01, 0.03, 0.98, 0.93])
    ax.axis("off")

    run_count = max(1, len(selected_runs))
    header_wrap_width = 12
    value_wrap_width = 15
    row_label_wrap_width = 28
    first_col_width = 0.24 if run_count <= 4 else 0.26
    remaining_width = 1.0 - first_col_width
    col_widths = [first_col_width] + [remaining_width / run_count] * run_count

    column_labels = ["Hyperparameter"] + [
        wrap_label(run["name"], header_wrap_width) for run in selected_runs
    ]
    cell_text: list[list[str]] = []
    for hyperparameter_name in hyperparameter_names:
        row_label = hyperparameter_labels.get(hyperparameter_name, hyperparameter_name)
        row = [wrap_label(row_label, row_label_wrap_width)]
        for run in selected_runs:
            row.append(
                wrap_label(
                    hyperparameter_value_for_run(run, hyperparameter_name),
                    value_wrap_width,
                )
            )
        cell_text.append(row)

    table = ax.table(
        cellText=cell_text,
        colLabels=column_labels,
        colWidths=col_widths,
        cellLoc="center",
        bbox=[0.005, 0.055, 0.99, 0.89],
    )
    table.auto_set_font_size(False)
    font_size = 10.5 if run_count <= 4 else 10 if run_count <= 6 else 9
    table.set_fontsize(font_size)
    row_scale = 1.8 if len(hyperparameter_names) <= 10 else 1.65
    table.scale(1, row_scale)

    for col_index in range(len(column_labels)):
        header_cell = table[(0, col_index)]
        header_cell.set_facecolor("#e9eef5")
        header_cell.get_text().set_weight("bold")
        if col_index != 0:
            header_cell.get_text().set_linespacing(0.85)

    for row_index in range(1, len(cell_text) + 1):
        first_col_cell = table[(row_index, 0)]
        first_col_cell.set_facecolor("#f7f7f7")
        first_col_cell.get_text().set_weight("bold")
        first_col_cell.get_text().set_ha("left")
        first_col_cell.get_text().set_x(0.02)
        for col_index in range(1, len(column_labels)):
            table[(row_index, col_index)].get_text().set_linespacing(0.85)

    add_page_number(fig, page_number)
    writer.save_figure(fig, dpi=dpi)
    plt.close(fig)
    return page_number + 1


def render_image_pages(
    writer: "ReportWriter",
    selected_runs: list[dict[str, Any]],
    image_keys: list[str],
    max_rows: int,
    max_columns: int,
    run_name_max_chars_per_line: int,
    image_max_side: int,
    image_dpi_percent: float | None,
    per_image_metric: dict[str, Any],
    dpi: int,
    page_number: int,
    page_title: str = "Comet Run Image Comparison",
    face_bbox_map: dict[str, list[float]] | None = None,
    face_closeups: dict[str, Any] | None = None,
) -> int:
    if not image_keys:
        fig = plt.figure(figsize=LANDSCAPE_A4)
        fig.text(0.5, 0.55, "No downloaded images found for the selected runs.", ha="center", va="center", fontsize=18)
        fig.text(0.5, 0.50, "Only metric charts will be included after this page.", ha="center", va="center", fontsize=12)
        add_step_footnote(fig, selected_runs)
        add_page_number(fig, page_number)
        writer.save_figure(fig, dpi=dpi)
        plt.close(fig)
        return page_number + 1

    run_chunks = chunked(selected_runs, max_rows)
    image_chunks = chunked(image_keys, max_columns)

    for image_chunk in image_chunks:
        for run_chunk in run_chunks:
            row_count = len(run_chunk)
            col_count = len(image_chunk)
            fig = plt.figure(figsize=LANDSCAPE_A4)
            grid = GridSpec(
                row_count + 1,
                col_count + 1,
                figure=fig,
                width_ratios=[0.22] + [1.0] * col_count,
                height_ratios=[0.22] + [1.0] * row_count,
                left=0.04,
                right=0.985,
                top=0.94,
                bottom=0.10,
                wspace=0.04,
                hspace=0.06,
            )

            fig.text(
                0.04,
                0.965,
                page_title,
                ha="left",
                va="top",
                fontsize=13,
                weight="bold",
            )

            corner_ax = fig.add_subplot(grid[0, 0])
            corner_ax.axis("off")

            for col_index, image_key in enumerate(image_chunk, start=1):
                header_ax = fig.add_subplot(grid[0, col_index])
                header_ax.axis("off")
                header_ax.text(
                    0.5,
                    0.5,
                    wrap_label(image_key, 18),
                    ha="center",
                    va="center",
                    fontsize=10,
                    weight="bold",
                )

            for row_index, run in enumerate(run_chunk, start=1):
                run_label_ax = fig.add_subplot(grid[row_index, 0])
                run_label_ax.axis("off")
                run_label_ax.text(
                    0.5,
                    0.5,
                    wrap_label(run["name"], run_name_max_chars_per_line),
                    ha="center",
                    va="center",
                    fontsize=11,
                    weight="bold",
                    rotation=90,
                    linespacing=0.8,
                )

                for col_index, image_key in enumerate(image_chunk, start=1):
                    cell_ax = fig.add_subplot(grid[row_index, col_index])
                    image_meta = run["image_map"].get(image_key)
                    if image_meta is None or not image_meta["path"].is_file():
                        render_missing_cell(cell_ax)
                        continue
                    image = load_image_for_report(
                        image_meta["path"],
                        image_max_side,
                        image_dpi_percent,
                        crop_bbox=(
                            face_bbox_map[image_key]
                            if face_bbox_map is not None
                            else None
                        ),
                        crop_padding_ratio=(
                            float(face_closeups["padding_ratio"])
                            if face_closeups is not None
                            else 0.0
                        ),
                        bbox_coordinate_size=(
                            face_closeups.get("bbox_coordinate_size")
                            if face_closeups is not None
                            else None
                        ),
                    )
                    cell_ax.imshow(image)
                    cell_ax.set_xticks([])
                    cell_ax.set_yticks([])
                    cell_ax.set_frame_on(False)
                    if per_image_metric.get("enabled", False):
                        metric_map = run.get("per_image_metric_map", {})
                        if image_key not in metric_map:
                            raise ConfigError(
                                f"{run['name']}: missing image metric for {image_key}"
                            )
                        decimals = int(per_image_metric.get("decimals", 3))
                        label = str(per_image_metric.get("label") or "ID")
                        cell_ax.text(
                            0.975,
                            0.975,
                            f"{label} {metric_map[image_key]:.{decimals}f}",
                            transform=cell_ax.transAxes,
                            ha="right",
                            va="top",
                            fontsize=8.2,
                            color="white",
                            weight="bold",
                            bbox={
                                "boxstyle": "round,pad=0.25",
                                "facecolor": "#111827",
                                "edgecolor": "white",
                                "linewidth": 0.55,
                                "alpha": 0.90,
                            },
                        )

            add_step_footnote(fig, run_chunk)
            add_page_number(fig, page_number)
            writer.save_figure(fig, dpi=dpi)
            plt.close(fig)
            page_number += 1

    return page_number


def render_metric_pages(
    writer: "ReportWriter",
    selected_runs: list[dict[str, Any]],
    metrics: list[str],
    metric_point_labels: dict[str, Any],
    dpi: int,
    page_number: int,
) -> int:
    metric_pages = chunked(metrics, MAX_CHARTS_PER_PAGE) if metrics else [[]]

    for metric_page in metric_pages:
        metric_count = len(metric_page)
        if metric_count == 0:
            fig = plt.figure(figsize=LANDSCAPE_A4)
            fig.text(0.5, 0.56, "No numeric metrics were available for charting.", ha="center", va="center", fontsize=18)
            fig.text(0.5, 0.50, "Set 'key_metrics' in the PDF config to request specific metrics.", ha="center", va="center", fontsize=11)
            add_page_number(fig, page_number)
            writer.save_figure(fig, dpi=dpi)
            plt.close(fig)
            page_number += 1
            continue

        chart_cols = 2 if metric_count > 1 else 1
        chart_rows = math.ceil(metric_count / chart_cols)
        fig, axes = plt.subplots(chart_rows, chart_cols, figsize=LANDSCAPE_A4)
        axes_list = axes.flatten().tolist() if hasattr(axes, "flatten") else [axes]
        fig.subplots_adjust(left=0.07, right=0.985, top=0.91, bottom=0.08, wspace=0.25, hspace=0.38)
        fig.suptitle("Comet Metric Comparison", fontsize=14, weight="bold", y=0.965)

        for ax, metric_name in zip(axes_list, metric_page):
            plotted_any = False
            for run_index, run in enumerate(selected_runs):
                series = extract_metric_series(run["export_run"].get("metrics", {}).get(metric_name))
                if not series:
                    continue
                steps, values = zip(*series)
                line, = ax.plot(
                    steps,
                    values,
                    linewidth=1.8,
                    marker="o",
                    markersize=4,
                    label=run["name"],
                )
                if metric_point_labels.get("enabled", False):
                    decimals = int(metric_point_labels.get("decimals", 3))
                    color = line.get_color()
                    last_index = len(series) - 1
                    max_index = max(range(len(series)), key=lambda index: values[index])
                    label_indices = [(last_index, "last")]
                    if max_index != last_index:
                        label_indices.append((max_index, "max"))
                    vertical_offsets = [-15, 12, -22, 18]
                    for point_order, (point_index, point_kind) in enumerate(label_indices):
                        y_offset = vertical_offsets[
                            (run_index + point_order) % len(vertical_offsets)
                        ]
                        ax.annotate(
                            f"{point_kind} {values[point_index]:.{decimals}f}",
                            xy=(steps[point_index], values[point_index]),
                            xytext=(7, y_offset),
                            textcoords="offset points",
                            ha="left",
                            va="center",
                            fontsize=7.5,
                            color="white",
                            weight="bold",
                            bbox={
                                "boxstyle": "round,pad=0.22",
                                "facecolor": color,
                                "edgecolor": "white",
                                "linewidth": 0.5,
                                "alpha": 0.92,
                            },
                            arrowprops={
                                "arrowstyle": "-",
                                "color": color,
                                "linewidth": 0.7,
                            },
                            clip_on=False,
                        )
                plotted_any = True

            if not plotted_any:
                ax.text(0.5, 0.5, "No numeric series", ha="center", va="center", fontsize=11)
                ax.set_axis_off()
                continue

            ax.set_title(metric_name, fontsize=11, weight="bold")
            ax.set_xlabel("Step")
            ax.set_ylabel("Value")
            ax.grid(True, alpha=0.3)
            ax.legend(fontsize=8)
            ax.margins(x=0.08, y=0.12)

        for ax in axes_list[metric_count:]:
            ax.set_axis_off()

        add_page_number(fig, page_number)
        writer.save_figure(fig, dpi=dpi)
        plt.close(fig)
        page_number += 1

    return page_number


def grouped_metric_means(
    run: dict[str, Any],
    metric_column: str,
    group_kind: str,
) -> dict[str, float]:
    grouped: dict[str, list[float]] = {}
    for row in run.get("per_image_metric_rows", []):
        if group_kind == "identity":
            group = str(row.get("identity") or "Unknown").strip().title()
        elif group_kind == "prompt":
            group = prompt_group_label(row)
        else:
            raise ConfigError(f"Unknown metric grouping: {group_kind}")
        try:
            value = float(row[metric_column])
        except (TypeError, ValueError, KeyError) as exc:
            raise ConfigError(
                f"{run['name']}: invalid {metric_column!r} in grouped table"
            ) from exc
        grouped.setdefault(group, []).append(value)
    return {group: sum(values) / len(values) for group, values in grouped.items()}


def render_group_average_page(
    writer: "ReportWriter",
    selected_runs: list[dict[str, Any]],
    table_config: dict[str, Any],
    dpi: int,
    page_number: int,
) -> int:
    if not table_config.get("enabled", False):
        return page_number

    metric_column = str(table_config.get("metric_column") or "id_sim")
    metric_label = str(
        table_config.get("metric_label") or "Subject-v2 ID similarity"
    )
    decimals = int(table_config.get("decimals", 3))

    identity_values = [
        grouped_metric_means(run, metric_column, "identity")
        for run in selected_runs
    ]
    prompt_values = [
        grouped_metric_means(run, metric_column, "prompt")
        for run in selected_runs
    ]
    identities = sorted(
        set().union(*(values.keys() for values in identity_values)),
        key=natural_sort_key,
    )
    prompts = sorted(
        set().union(*(values.keys() for values in prompt_values)),
        key=natural_sort_key,
    )

    fig = plt.figure(figsize=LANDSCAPE_A4)
    fig.text(
        0.02,
        0.982,
        f"Mean {metric_label} by Identity and Prompt",
        ha="left",
        va="top",
        fontsize=14,
        weight="bold",
    )
    axes = [
        fig.add_axes([0.02, 0.10, 0.40, 0.80]),
        fig.add_axes([0.44, 0.10, 0.54, 0.80]),
    ]
    run_labels = [wrap_label(run["name"], 10) for run in selected_runs]

    def draw_table(
        ax: plt.Axes,
        title: str,
        groups: list[str],
        values_by_run: list[dict[str, float]],
        first_col_width: float,
    ) -> None:
        ax.axis("off")
        ax.set_title(title, fontsize=12, weight="bold", pad=8)
        rows: list[list[str]] = []
        for group in groups:
            rows.append(
                [wrap_label(group, 18)]
                + [
                    (
                        f"{values[group]:.{decimals}f}"
                        if group in values
                        else "n/a"
                    )
                    for values in values_by_run
                ]
            )
        remaining = 1.0 - first_col_width
        table = ax.table(
            cellText=rows,
            colLabels=[title[:-1] if title.endswith("s") else title] + run_labels,
            colWidths=[first_col_width]
            + [remaining / len(selected_runs)] * len(selected_runs),
            cellLoc="center",
            bbox=[0.0, 0.0, 1.0, 0.96],
        )
        table.auto_set_font_size(False)
        table.set_fontsize(8.5)
        table.scale(1, 1.35)
        for col in range(1 + len(selected_runs)):
            table[(0, col)].set_facecolor("#e9eef5")
            table[(0, col)].get_text().set_weight("bold")
        for row_index, group in enumerate(groups, start=1):
            table[(row_index, 0)].set_facecolor("#f7f7f7")
            table[(row_index, 0)].get_text().set_weight("bold")
            numeric = [values.get(group) for values in values_by_run]
            finite = [value for value in numeric if value is not None]
            if finite:
                row_max = max(finite)
                for run_index, value in enumerate(numeric, start=1):
                    if value is not None and math.isclose(value, row_max):
                        table[(row_index, run_index)].set_facecolor("#e3f3e8")
                        table[(row_index, run_index)].get_text().set_weight("bold")

    draw_table(axes[0], "Identities", identities, identity_values, 0.27)
    draw_table(axes[1], "Prompts", prompts, prompt_values, 0.26)
    add_step_footnote(fig, selected_runs)
    add_page_number(fig, page_number)
    writer.save_figure(fig, dpi=dpi)
    plt.close(fig)
    return page_number + 1


class ReportWriter:
    def save_figure(self, fig: plt.Figure, dpi: int) -> None:
        raise NotImplementedError

    def close(self) -> None:
        raise NotImplementedError


class PdfReportWriter(ReportWriter):
    def __init__(self, output_path: Path) -> None:
        self.output_path = output_path
        self._pdf = PdfPages(output_path)

    def save_figure(self, fig: plt.Figure, dpi: int) -> None:
        self._pdf.savefig(fig, dpi=dpi)

    def close(self) -> None:
        self._pdf.close()


class PptxReportWriter(ReportWriter):
    def __init__(self, output_path: Path) -> None:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise ConfigError(
                "PPTX output requires python-pptx. Install it with 'pip install python-pptx'."
            ) from exc

        self.output_path = output_path
        self._tmpdir = tempfile.TemporaryDirectory(prefix="comet_report_pptx_")
        self._slide_index = 0
        self._presentation = Presentation()
        self._presentation.slide_width = int(LANDSCAPE_A4[0] * 914400)
        self._presentation.slide_height = int(LANDSCAPE_A4[1] * 914400)
        self._blank_layout = self._presentation.slide_layouts[6]

    def save_figure(self, fig: plt.Figure, dpi: int) -> None:
        slide_image = Path(self._tmpdir.name) / f"slide_{self._slide_index:04d}.png"
        fig.savefig(slide_image, dpi=dpi)
        slide = self._presentation.slides.add_slide(self._blank_layout)
        slide.shapes.add_picture(
            str(slide_image),
            0,
            0,
            width=self._presentation.slide_width,
            height=self._presentation.slide_height,
        )
        self._slide_index += 1

    def close(self) -> None:
        self._presentation.save(self.output_path)
        self._tmpdir.cleanup()


def make_report_writer(output_path: Path, output_format: str) -> ReportWriter:
    if output_format == "pdf":
        return PdfReportWriter(output_path)
    if output_format == "pptx":
        return PptxReportWriter(output_path)
    raise ConfigError(f"Unsupported output format: {output_format}")


def ensure_output_suffix(output_path: Path, output_format: str) -> Path:
    expected_suffix = ".pptx" if output_format == "pptx" else ".pdf"
    if output_path.suffix.lower() == expected_suffix:
        return output_path
    return output_path.with_suffix(expected_suffix)


def resolve_image_dpi_percent(
    config_value: float | None,
    override_value: float | None,
) -> float | None:
    if override_value is None:
        return config_value
    return parse_optional_positive_number(override_value, "image_dpi_percent override")


def build_report(
    config_path: Path,
    output_path: Path,
    dpi: int,
    image_max_side: int,
    image_dpi_percent_override: float | None = None,
) -> Path:
    config, _ = load_pdf_config(config_path)
    export_json_path = Path(config["export_json"]).resolve()
    export_payload = load_json(export_json_path)
    if not isinstance(export_payload, dict):
        raise ConfigError(f"Export JSON must be an object: {export_json_path}")

    selected_runs, image_root_dir = prepare_selected_runs(
        config["runs"],
        export_payload,
        export_json_path,
        config["ignore_mask"],
    )
    attach_per_image_metric_data(
        selected_runs,
        image_root_dir,
        config["per_image_metric"],
    )
    image_keys = collect_image_keys(selected_runs)
    face_bbox_map = load_face_bbox_map(config["face_closeups"], image_keys)
    metrics = select_numeric_metrics(selected_runs, config["key_metrics"])
    image_dpi_percent = resolve_image_dpi_percent(
        config["image_dpi_percent"],
        image_dpi_percent_override,
    )

    output_path = ensure_output_suffix(output_path.expanduser().resolve(), config["output_format"])
    output_path.parent.mkdir(parents=True, exist_ok=True)

    page_number = 1
    writer = make_report_writer(output_path, config["output_format"])
    try:
        page_number = render_image_pages(
            writer=writer,
            selected_runs=selected_runs,
            image_keys=image_keys,
            max_rows=config["max_rows"],
            max_columns=config["max_columns"],
            run_name_max_chars_per_line=config["run_name_max_chars_per_line"],
            image_max_side=image_max_side,
            image_dpi_percent=image_dpi_percent,
            per_image_metric=config["per_image_metric"],
            dpi=dpi,
            page_number=page_number,
        )
        if config["face_closeups"]["enabled"]:
            page_number = render_image_pages(
                writer=writer,
                selected_runs=selected_runs,
                image_keys=image_keys,
                max_rows=config["max_rows"],
                max_columns=config["max_columns"],
                run_name_max_chars_per_line=config["run_name_max_chars_per_line"],
                image_max_side=image_max_side,
                image_dpi_percent=image_dpi_percent,
                per_image_metric=config["per_image_metric"],
                dpi=dpi,
                page_number=page_number,
                page_title=str(config["face_closeups"]["title"]),
                face_bbox_map=face_bbox_map,
                face_closeups=config["face_closeups"],
            )
        page_number = render_markdown_pages(
            writer=writer,
            markdown_source=config["markdown_source"],
            dpi=dpi,
            page_number=page_number,
        )
        page_number = render_metric_pages(
            writer=writer,
            selected_runs=selected_runs,
            metrics=metrics,
            metric_point_labels=config["metric_point_labels"],
            dpi=dpi,
            page_number=page_number,
        )
        page_number = render_group_average_page(
            writer=writer,
            selected_runs=selected_runs,
            table_config=config["group_average_tables"],
            dpi=dpi,
            page_number=page_number,
        )
        render_hyperparameter_page(
            writer=writer,
            selected_runs=selected_runs,
            hyperparameter_names=config["key_hyperparameters"],
            hyperparameter_labels=config["hyperparameter_labels"],
            run_name_max_chars_per_line=config["run_name_max_chars_per_line"],
            dpi=dpi,
            page_number=page_number,
        )
    finally:
        writer.close()

    return output_path


def main() -> int:
    args = parse_args()
    try:
        output_path = args.output if args.output is not None else args.output_pdf
        saved_report = build_report(
            config_path=args.config.resolve(),
            output_path=output_path,
            dpi=args.dpi,
            image_max_side=args.image_max_side,
            image_dpi_percent_override=args.image_dpi_percent,
        )
    except ConfigError as exc:
        print(f"Config error: {exc}")
        return 2
    except Exception as exc:
        print(f"Failed to build report: {exc}")
        return 1

    print(f"Saved report to {saved_report}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
