#!/usr/bin/env python3
"""
Build a landscape report from the output of export_comet_runs.py.

Usage:
    python diffusion_template/comet_utils/build_comet_report_pdf.py \
        --config diffusion_template/comet_utils/comet_pdf_config_template.json \
        --output diffusion_template/comet_data/comet_report.pdf

If the config omits "runs" or sets it to an empty list, all runs from the
export JSON are included in the report in export order.
"""

from __future__ import annotations

import argparse
import json
import math
import re
import tempfile
import textwrap
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.gridspec import GridSpec
from PIL import Image, ImageOps

SCRIPT_DIR = Path(__file__).resolve().parent
DEFAULT_CONFIG_PATH = SCRIPT_DIR / "comet_pdf_config_template.json"
DEFAULT_EXPORT_JSON = SCRIPT_DIR.parent / "comet_data" / "comet_runs_export.json"
LANDSCAPE_A4 = (11.69, 8.27)
IMG_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
PLACEHOLDER_RUN_ID = "REPLACE_WITH_COMET_EXPERIMENT_KEY"
MAX_CHARTS_PER_PAGE = 6
PREFERRED_METRICS = [
    "train_loss",
    "loss",
    "val_loss",
    "train/loss",
    "val/loss",
    "accuracy",
    "val_accuracy",
    "train_accuracy",
    "general/steps_per_sec",
]
DEFAULT_KEY_HYPERPARAMETERS = [
    "step_shown",
    "model.weight_dtype",
    "model.rank",
    "dataloaders.train.batch_size",
    "branched_attn_weight_mode",
    "branched_attn_new_weight_kind",
    "train_branched_ca_lora",
    "lr_for_lora",
    "loss_kind",
    "lambda_face",
    "strict_face_routing",
    "trainer.masked_loss_step",
]


class ConfigError(ValueError):
    """Raised when the PDF configuration is invalid."""


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create a landscape report comparing exported Comet images and metrics."
    )
    parser.add_argument(
        "--config",
        type=Path,
        required=True,
        help=f"Path to the JSON config file. Template: {DEFAULT_CONFIG_PATH}",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="Path to the report file that will be created. Supports .pdf and .pptx.",
    )
    parser.add_argument(
        "--output-pdf",
        type=Path,
        default=None,
        help="Deprecated alias for --output.",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=200,
        help="Render DPI for report pages. Default: 200",
    )
    parser.add_argument(
        "--image-max-side",
        type=int,
        default=768,
        help="Maximum side length used when loading images. Default: 768",
    )
    parser.add_argument(
        "--image-dpi-percent",
        type=float,
        default=None,
        help=(
            "Optional override for the JSON image_dpi_percent setting. "
            "Examples: 100 preserves original size, 50 uses half size."
        ),
    )
    args = parser.parse_args()
    if args.output is None and args.output_pdf is None:
        parser.error("one of --output or --output-pdf is required")
    if args.output is not None and args.output_pdf is not None:
        parser.error("use only one of --output or --output-pdf")
    return args


def load_json(path: Path) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError as exc:
        raise ConfigError(f"JSON file does not exist: {path}") from exc
    except json.JSONDecodeError as exc:
        raise ConfigError(f"Invalid JSON file: {path}") from exc


def resolve_path(base_dir: Path, raw_path: str | None, fallback: Path | None = None) -> Path:
    if raw_path:
        path = Path(raw_path).expanduser()
        if not path.is_absolute():
            path = (base_dir / path).resolve()
        return path
    if fallback is None:
        raise ConfigError("A required path value is missing.")
    return fallback.resolve()


def parse_positive_int(value: Any, field_name: str) -> int:
    if isinstance(value, bool):
        raise ConfigError(f"{field_name} must be a positive integer")
    try:
        parsed = int(value)
    except (TypeError, ValueError) as exc:
        raise ConfigError(f"{field_name} must be a positive integer") from exc
    if parsed <= 0:
        raise ConfigError(f"{field_name} must be a positive integer")
    return parsed


def parse_bool(value: Any, field_name: str) -> bool:
    if isinstance(value, bool):
        return value
    raise ConfigError(f"{field_name} must be true or false")


def parse_optional_positive_number(value: Any, field_name: str) -> float | None:
    if value is None:
        return None
    if isinstance(value, bool):
        raise ConfigError(f"{field_name} must be a positive number")
    try:
        parsed = float(value)
    except (TypeError, ValueError) as exc:
        raise ConfigError(f"{field_name} must be a positive number") from exc
    if parsed <= 0:
        raise ConfigError(f"{field_name} must be a positive number")
    return parsed


def parse_output_format(value: Any) -> str:
    if value is None:
        return "pdf"
    text = str(value).strip().lower()
    if text in {"pdf", "ppt", "pptx"}:
        return "pptx" if text in {"ppt", "pptx"} else "pdf"
    raise ConfigError("output_format must be 'pdf' or 'pptx'")


def natural_sort_key(text: str) -> list[Any]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", text)]


def chunked(items: list[Any], size: int) -> list[list[Any]]:
    return [items[index : index + size] for index in range(0, len(items), size)]


def wrap_label(text: str, width: int) -> str:
    normalized = " ".join(str(text).split())
    return "\n".join(textwrap.wrap(normalized, width=width)[:4]) if normalized else ""


def load_image_for_report(
    path: Path,
    max_side: int,
    image_dpi_percent: float | None,
) -> Image.Image:
    image = Image.open(path).convert("RGB")
    if image_dpi_percent is not None:
        scale = image_dpi_percent / 100.0
        new_size = (
            max(1, round(image.size[0] * scale)),
            max(1, round(image.size[1] * scale)),
        )
        if new_size != image.size:
            image = image.resize(new_size, Image.LANCZOS)
    elif max(image.size) > max_side:
        image.thumbnail((max_side, max_side), Image.LANCZOS)
    return image


def placeholder_image(size: tuple[int, int] = (512, 512)) -> Image.Image:
    image = Image.new("RGB", size, "#f2f2f2")
    return ImageOps.expand(image, border=2, fill="#cccccc")


def display_label_for_image(file_name: str) -> str:
    path = Path(file_name)
    label = path.stem if path.stem else path.name
    label = re.sub(r"\s*\(\d+\)$", "", label).strip()
    label = re.sub(r"_{1,}\d+$", "", label).strip()
    return label or path.stem or path.name


def canonical_image_key(file_name: str) -> str:
    return display_label_for_image(file_name)


def is_mask_image_name(file_name: str) -> bool:
    raw_stem = Path(file_name).stem
    normalized_stem = display_label_for_image(file_name)
    return raw_stem.endswith("_mask") or normalized_stem.endswith("_mask")


def run_display_name(run_cfg: dict[str, Any], export_run: dict[str, Any]) -> str:
    return str(run_cfg.get("run_name") or export_run.get("name") or export_run.get("id"))


def normalize_string_list(value: Any, field_name: str) -> list[str]:
    if value is None:
        return []
    if not isinstance(value, list):
        raise ConfigError(f"'{field_name}' must be a list when provided.")
    return [str(item).strip() for item in value if str(item).strip()]


def select_numeric_metrics(
    selected_runs: list[dict[str, Any]],
    configured_metrics: list[str] | None,
) -> list[str]:
    if configured_metrics:
        return [metric for metric in configured_metrics if metric]

    run_metric_sets = []
    for run in selected_runs:
        numeric_metrics: set[str] = set()
        for metric_name, entries in run["export_run"].get("metrics", {}).items():
            if extract_metric_series(entries):
                numeric_metrics.add(metric_name)
        run_metric_sets.append(numeric_metrics)

    if not run_metric_sets:
        return []

    common_metrics = set.intersection(*run_metric_sets) if run_metric_sets else set()
    candidate_metrics = common_metrics or set.union(*run_metric_sets)

    def metric_priority(metric_name: str) -> tuple[int, str]:
        if metric_name in PREFERRED_METRICS:
            return (PREFERRED_METRICS.index(metric_name), metric_name.lower())
        return (len(PREFERRED_METRICS) + 1, metric_name.lower())

    return sorted(candidate_metrics, key=metric_priority)[:MAX_CHARTS_PER_PAGE]


def extract_metric_series(entries: Any) -> list[tuple[float, float]]:
    series: list[tuple[float, float]] = []
    if not isinstance(entries, list):
        return series
    for entry in entries:
        if not isinstance(entry, dict):
            continue
        step = entry.get("step")
        value = entry.get("value")
        if isinstance(step, bool) or isinstance(value, bool):
            continue
        try:
            step_num = float(step)
            value_num = float(value)
        except (TypeError, ValueError):
            continue
        series.append((step_num, value_num))
    series.sort(key=lambda pair: pair[0])
    return series


def load_pdf_config(path: Path) -> tuple[dict[str, Any], Path]:
    raw_config = load_json(path)
    if not isinstance(raw_config, dict):
        raise ConfigError("PDF config must be a JSON object.")

    config_dir = path.resolve().parent
    runs = raw_config.get("runs", [])
    if runs is None:
        runs = []
    if not isinstance(runs, list):
        raise ConfigError("'runs' must be a list when provided.")

    normalized_runs: list[dict[str, Any]] = []
    for index, run in enumerate(runs, start=1):
        if not isinstance(run, dict):
            raise ConfigError(f"runs[{index}] must be an object.")
        run_id = str(run.get("run_id", "")).strip()
        if not run_id:
            raise ConfigError(f"runs[{index}].run_id is required.")
        if run_id == PLACEHOLDER_RUN_ID:
            raise ConfigError(
                f"runs[{index}].run_id still uses the template placeholder {PLACEHOLDER_RUN_ID}."
            )
        run_name = run.get("run_name")
        normalized_runs.append(
            {
                "run_id": run_id,
                "run_name": None if run_name in (None, "") else str(run_name).strip(),
            }
        )

    configured_metrics = normalize_string_list(raw_config.get("key_metrics"), "key_metrics")
    if not configured_metrics:
        configured_metrics = None
    configured_hyperparameters = normalize_string_list(
        raw_config.get("key_hyperparameters", DEFAULT_KEY_HYPERPARAMETERS),
        "key_hyperparameters",
    )

    normalized_config = {
        "export_json": resolve_path(config_dir, raw_config.get("export_json"), DEFAULT_EXPORT_JSON),
        "output_format": parse_output_format(raw_config.get("output_format", "pdf")),
        "max_columns": parse_positive_int(raw_config.get("max_columns", 4), "max_columns"),
        "max_rows": parse_positive_int(raw_config.get("max_rows", 3), "max_rows"),
        "ignore_mask": parse_bool(raw_config.get("ignore_mask", True), "ignore_mask"),
        "image_dpi_percent": parse_optional_positive_number(
            raw_config.get("image_dpi_percent"),
            "image_dpi_percent",
        ),
        "run_name_max_chars_per_line": parse_positive_int(
            raw_config.get("run_name_max_chars_per_line", 16),
            "run_name_max_chars_per_line",
        ),
        "key_metrics": configured_metrics,
        "key_hyperparameters": configured_hyperparameters,
        "runs": normalized_runs,
    }
    return normalized_config, config_dir


def prepare_selected_runs(
    config_runs: list[dict[str, Any]],
    export_payload: dict[str, Any],
    export_json_path: Path,
    ignore_mask: bool,
) -> tuple[list[dict[str, Any]], Path]:
    export_runs = export_payload.get("runs")
    if not isinstance(export_runs, list):
        raise ConfigError(f"'runs' is missing or invalid in export JSON: {export_json_path}")

    ordered_export_runs: list[dict[str, Any]] = []
    run_by_id: dict[str, dict[str, Any]] = {}
    for run in export_runs:
        if isinstance(run, dict) and run.get("id"):
            run_id = str(run["id"])
            ordered_export_runs.append(run)
            run_by_id[run_id] = run

    images_root = export_payload.get("output_dir")
    if not images_root:
        raise ConfigError(f"'output_dir' is missing in export JSON: {export_json_path}")
    image_root_dir = Path(images_root).expanduser().resolve()

    selected_runs: list[dict[str, Any]] = []
    if not config_runs:
        for export_run in ordered_export_runs:
            selected_runs.append(
                {
                    "id": export_run["id"],
                    "name": str(export_run.get("name") or export_run.get("id")),
                    "export_run": export_run,
                    "image_map": build_run_image_map(export_run, image_root_dir, ignore_mask),
                }
            )
        if not selected_runs:
            raise ConfigError(f"No runs were found in export JSON: {export_json_path}")
        return selected_runs, image_root_dir

    missing_ids: list[str] = []
    for run_cfg in config_runs:
        export_run = run_by_id.get(run_cfg["run_id"])
        if export_run is None:
            missing_ids.append(run_cfg["run_id"])
            continue
        selected_runs.append(
            {
                "id": export_run["id"],
                "name": run_display_name(run_cfg, export_run),
                "export_run": export_run,
                "image_map": build_run_image_map(export_run, image_root_dir, ignore_mask),
            }
        )

    if missing_ids:
        raise ConfigError(f"Run IDs not found in export JSON: {', '.join(missing_ids)}")
    if not selected_runs:
        raise ConfigError("No runs were selected for the PDF report.")

    return selected_runs, image_root_dir


def build_run_image_map(
    export_run: dict[str, Any],
    image_root_dir: Path,
    ignore_mask: bool,
) -> dict[str, dict[str, Any]]:
    image_map: dict[str, dict[str, Any]] = {}
    for image_info in export_run.get("downloaded_images", []):
        if not isinstance(image_info, dict):
            continue
        file_name = str(image_info.get("file_name") or "").strip()
        if not file_name:
            continue
        if ignore_mask and is_mask_image_name(file_name):
            continue
        key = canonical_image_key(file_name)
        relative_path = image_info.get("saved_path")
        if relative_path:
            image_path = (image_root_dir / relative_path).resolve()
        else:
            output_folder = export_run.get("output_folder", "")
            image_path = (image_root_dir / str(output_folder) / file_name).resolve()
        image_map[key] = {
            "file_name": file_name,
            "display_name": display_label_for_image(file_name),
            "path": image_path,
        }
    return image_map


def collect_image_keys(selected_runs: list[dict[str, Any]]) -> list[str]:
    image_keys: set[str] = set()
    for run in selected_runs:
        image_keys.update(run["image_map"].keys())
    return sorted(image_keys, key=natural_sort_key)


def add_page_number(fig: plt.Figure, page_number: int) -> None:
    fig.text(0.5, 0.015, f"Page {page_number}", ha="center", va="bottom", fontsize=9)


def step_display_value_for_run(run: dict[str, Any]) -> str:
    export_run = run["export_run"]
    selection = export_run.get("step_selection")
    requested = export_run.get("requested_step_number")
    resolved = export_run.get("resolved_step_number")

    if isinstance(selection, dict):
        requested = selection.get("requested_step_number", requested)
        resolved = selection.get("resolved_step_number", resolved)

    if resolved is None and requested is None:
        return "step unavailable"
    if resolved is None:
        return f"no image step for requested {requested}"
    if requested is None or requested == resolved:
        return f"step {resolved}"
    return f"step {resolved} (requested {requested})"


def step_note_for_run(run: dict[str, Any]) -> str:
    return f"{run['name']}: {step_display_value_for_run(run)}"


def add_step_footnote(fig: plt.Figure, runs: list[dict[str, Any]]) -> None:
    if not runs:
        return
    note = "Step notes: " + "; ".join(step_note_for_run(run) for run in runs)
    wrapped = "\n".join(textwrap.wrap(note, width=135))
    fig.text(0.04, 0.032, wrapped, ha="left", va="bottom", fontsize=8, color="#444444")


def render_missing_cell(ax: plt.Axes) -> None:
    ax.imshow(placeholder_image())
    ax.text(0.5, 0.5, "Missing", ha="center", va="center", fontsize=10, color="#555555")
    ax.set_xticks([])
    ax.set_yticks([])
    ax.set_frame_on(False)


def stringify_hyperparameter_value(value: Any) -> str:
    if value in (None, ""):
        return "n/a"
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        return f"{value:.6g}"
    if isinstance(value, (list, dict)):
        return json.dumps(value, ensure_ascii=True, sort_keys=True)
    text = str(value).strip()
    return text if text else "n/a"


def hyperparameter_value_for_run(run: dict[str, Any], hyperparameter_name: str) -> str:
    if hyperparameter_name == "step_shown":
        return step_display_value_for_run(run)

    export_run = run["export_run"]
    hyperparameters = export_run.get("hyperparameters", {})
    if not isinstance(hyperparameters, dict):
        return "n/a"

    return stringify_hyperparameter_value(hyperparameters.get(hyperparameter_name))


def render_hyperparameter_page(
    writer: "ReportWriter",
    selected_runs: list[dict[str, Any]],
    hyperparameter_names: list[str],
    run_name_max_chars_per_line: int,
    dpi: int,
    page_number: int,
) -> int:
    if not hyperparameter_names:
        return page_number

    fig = plt.figure(figsize=LANDSCAPE_A4)
    fig.text(0.02, 0.982, "Key Hyperparameters", ha="left", va="top", fontsize=14, weight="bold")
    ax = fig.add_axes([0.01, 0.03, 0.98, 0.93])
    ax.axis("off")

    run_count = max(1, len(selected_runs))
    header_wrap_width = 10
    value_wrap_width = 10
    row_label_wrap_width = 28
    first_col_width = 0.24 if run_count <= 4 else 0.26
    remaining_width = 1.0 - first_col_width
    col_widths = [first_col_width] + [remaining_width / run_count] * run_count

    column_labels = ["Hyperparameter"] + [
        wrap_label(run["name"], header_wrap_width) for run in selected_runs
    ]
    cell_text: list[list[str]] = []
    for hyperparameter_name in hyperparameter_names:
        row = [wrap_label(hyperparameter_name, row_label_wrap_width)]
        for run in selected_runs:
            row.append(
                wrap_label(
                    hyperparameter_value_for_run(run, hyperparameter_name),
                    value_wrap_width,
                )
            )
        cell_text.append(row)

    table = ax.table(
        cellText=cell_text,
        colLabels=column_labels,
        colWidths=col_widths,
        cellLoc="center",
        bbox=[0.005, 0.055, 0.99, 0.89],
    )
    table.auto_set_font_size(False)
    font_size = 12 if run_count <= 4 else 11 if run_count <= 6 else 10
    table.set_fontsize(font_size)
    row_scale = 1.8 if len(hyperparameter_names) <= 10 else 1.65
    table.scale(1, row_scale)

    for col_index in range(len(column_labels)):
        header_cell = table[(0, col_index)]
        header_cell.set_facecolor("#e9eef5")
        header_cell.get_text().set_weight("bold")
        if col_index != 0:
            header_cell.get_text().set_linespacing(0.85)

    for row_index in range(1, len(cell_text) + 1):
        first_col_cell = table[(row_index, 0)]
        first_col_cell.set_facecolor("#f7f7f7")
        first_col_cell.get_text().set_weight("bold")
        first_col_cell.get_text().set_ha("left")
        first_col_cell.get_text().set_x(0.02)
        for col_index in range(1, len(column_labels)):
            table[(row_index, col_index)].get_text().set_linespacing(0.85)

    add_page_number(fig, page_number)
    writer.save_figure(fig, dpi=dpi)
    plt.close(fig)
    return page_number + 1


def render_image_pages(
    writer: "ReportWriter",
    selected_runs: list[dict[str, Any]],
    image_keys: list[str],
    max_rows: int,
    max_columns: int,
    run_name_max_chars_per_line: int,
    image_max_side: int,
    image_dpi_percent: float | None,
    dpi: int,
    page_number: int,
) -> int:
    if not image_keys:
        fig = plt.figure(figsize=LANDSCAPE_A4)
        fig.text(0.5, 0.55, "No downloaded images found for the selected runs.", ha="center", va="center", fontsize=18)
        fig.text(0.5, 0.50, "Only metric charts will be included after this page.", ha="center", va="center", fontsize=12)
        add_step_footnote(fig, selected_runs)
        add_page_number(fig, page_number)
        writer.save_figure(fig, dpi=dpi)
        plt.close(fig)
        return page_number + 1

    run_chunks = chunked(selected_runs, max_rows)
    image_chunks = chunked(image_keys, max_columns)

    for image_chunk in image_chunks:
        for run_chunk in run_chunks:
            row_count = len(run_chunk)
            col_count = len(image_chunk)
            fig = plt.figure(figsize=LANDSCAPE_A4)
            grid = GridSpec(
                row_count + 1,
                col_count + 1,
                figure=fig,
                width_ratios=[0.22] + [1.0] * col_count,
                height_ratios=[0.22] + [1.0] * row_count,
                left=0.04,
                right=0.985,
                top=0.94,
                bottom=0.10,
                wspace=0.04,
                hspace=0.06,
            )

            fig.text(
                0.04,
                0.965,
                "Comet Run Image Comparison",
                ha="left",
                va="top",
                fontsize=13,
                weight="bold",
            )

            corner_ax = fig.add_subplot(grid[0, 0])
            corner_ax.axis("off")

            for col_index, image_key in enumerate(image_chunk, start=1):
                header_ax = fig.add_subplot(grid[0, col_index])
                header_ax.axis("off")
                header_ax.text(
                    0.5,
                    0.5,
                    wrap_label(image_key, 18),
                    ha="center",
                    va="center",
                    fontsize=10,
                    weight="bold",
                )

            for row_index, run in enumerate(run_chunk, start=1):
                run_label_ax = fig.add_subplot(grid[row_index, 0])
                run_label_ax.axis("off")
                run_label_ax.text(
                    0.5,
                    0.5,
                    wrap_label(run["name"], run_name_max_chars_per_line),
                    ha="center",
                    va="center",
                    fontsize=11,
                    weight="bold",
                    rotation=90,
                    linespacing=0.8,
                )

                for col_index, image_key in enumerate(image_chunk, start=1):
                    cell_ax = fig.add_subplot(grid[row_index, col_index])
                    image_meta = run["image_map"].get(image_key)
                    if image_meta is None or not image_meta["path"].is_file():
                        render_missing_cell(cell_ax)
                        continue
                    image = load_image_for_report(
                        image_meta["path"],
                        image_max_side,
                        image_dpi_percent,
                    )
                    cell_ax.imshow(image)
                    cell_ax.set_xticks([])
                    cell_ax.set_yticks([])
                    cell_ax.set_frame_on(False)

            add_step_footnote(fig, run_chunk)
            add_page_number(fig, page_number)
            writer.save_figure(fig, dpi=dpi)
            plt.close(fig)
            page_number += 1

    return page_number


def render_metric_pages(
    writer: "ReportWriter",
    selected_runs: list[dict[str, Any]],
    metrics: list[str],
    dpi: int,
    page_number: int,
) -> int:
    metric_pages = chunked(metrics, MAX_CHARTS_PER_PAGE) if metrics else [[]]

    for metric_page in metric_pages:
        metric_count = len(metric_page)
        if metric_count == 0:
            fig = plt.figure(figsize=LANDSCAPE_A4)
            fig.text(0.5, 0.56, "No numeric metrics were available for charting.", ha="center", va="center", fontsize=18)
            fig.text(0.5, 0.50, "Set 'key_metrics' in the PDF config to request specific metrics.", ha="center", va="center", fontsize=11)
            add_page_number(fig, page_number)
            writer.save_figure(fig, dpi=dpi)
            plt.close(fig)
            page_number += 1
            continue

        chart_cols = 2 if metric_count > 1 else 1
        chart_rows = math.ceil(metric_count / chart_cols)
        fig, axes = plt.subplots(chart_rows, chart_cols, figsize=LANDSCAPE_A4)
        axes_list = axes.flatten().tolist() if hasattr(axes, "flatten") else [axes]
        fig.subplots_adjust(left=0.07, right=0.985, top=0.91, bottom=0.08, wspace=0.25, hspace=0.38)
        fig.suptitle("Comet Metric Comparison", fontsize=14, weight="bold", y=0.965)

        for ax, metric_name in zip(axes_list, metric_page):
            plotted_any = False
            for run in selected_runs:
                series = extract_metric_series(run["export_run"].get("metrics", {}).get(metric_name))
                if not series:
                    continue
                steps, values = zip(*series)
                ax.plot(steps, values, linewidth=1.6, label=run["name"])
                plotted_any = True

            if not plotted_any:
                ax.text(0.5, 0.5, "No numeric series", ha="center", va="center", fontsize=11)
                ax.set_axis_off()
                continue

            ax.set_title(metric_name, fontsize=11, weight="bold")
            ax.set_xlabel("Step")
            ax.set_ylabel("Value")
            ax.grid(True, alpha=0.3)
            ax.legend(fontsize=8)

        for ax in axes_list[metric_count:]:
            ax.set_axis_off()

        add_page_number(fig, page_number)
        writer.save_figure(fig, dpi=dpi)
        plt.close(fig)
        page_number += 1

    return page_number


class ReportWriter:
    def save_figure(self, fig: plt.Figure, dpi: int) -> None:
        raise NotImplementedError

    def close(self) -> None:
        raise NotImplementedError


class PdfReportWriter(ReportWriter):
    def __init__(self, output_path: Path) -> None:
        self.output_path = output_path
        self._pdf = PdfPages(output_path)

    def save_figure(self, fig: plt.Figure, dpi: int) -> None:
        self._pdf.savefig(fig, dpi=dpi)

    def close(self) -> None:
        self._pdf.close()


class PptxReportWriter(ReportWriter):
    def __init__(self, output_path: Path) -> None:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise ConfigError(
                "PPTX output requires python-pptx. Install it with 'pip install python-pptx'."
            ) from exc

        self.output_path = output_path
        self._tmpdir = tempfile.TemporaryDirectory(prefix="comet_report_pptx_")
        self._slide_index = 0
        self._presentation = Presentation()
        self._presentation.slide_width = int(LANDSCAPE_A4[0] * 914400)
        self._presentation.slide_height = int(LANDSCAPE_A4[1] * 914400)
        self._blank_layout = self._presentation.slide_layouts[6]

    def save_figure(self, fig: plt.Figure, dpi: int) -> None:
        slide_image = Path(self._tmpdir.name) / f"slide_{self._slide_index:04d}.png"
        fig.savefig(slide_image, dpi=dpi)
        slide = self._presentation.slides.add_slide(self._blank_layout)
        slide.shapes.add_picture(
            str(slide_image),
            0,
            0,
            width=self._presentation.slide_width,
            height=self._presentation.slide_height,
        )
        self._slide_index += 1

    def close(self) -> None:
        self._presentation.save(self.output_path)
        self._tmpdir.cleanup()


def make_report_writer(output_path: Path, output_format: str) -> ReportWriter:
    if output_format == "pdf":
        return PdfReportWriter(output_path)
    if output_format == "pptx":
        return PptxReportWriter(output_path)
    raise ConfigError(f"Unsupported output format: {output_format}")


def ensure_output_suffix(output_path: Path, output_format: str) -> Path:
    expected_suffix = ".pptx" if output_format == "pptx" else ".pdf"
    if output_path.suffix.lower() == expected_suffix:
        return output_path
    return output_path.with_suffix(expected_suffix)


def resolve_image_dpi_percent(
    config_value: float | None,
    override_value: float | None,
) -> float | None:
    if override_value is None:
        return config_value
    return parse_optional_positive_number(override_value, "image_dpi_percent override")


def build_report(
    config_path: Path,
    output_path: Path,
    dpi: int,
    image_max_side: int,
    image_dpi_percent_override: float | None = None,
) -> Path:
    config, _ = load_pdf_config(config_path)
    export_json_path = Path(config["export_json"]).resolve()
    export_payload = load_json(export_json_path)
    if not isinstance(export_payload, dict):
        raise ConfigError(f"Export JSON must be an object: {export_json_path}")

    selected_runs, _ = prepare_selected_runs(
        config["runs"],
        export_payload,
        export_json_path,
        config["ignore_mask"],
    )
    image_keys = collect_image_keys(selected_runs)
    metrics = select_numeric_metrics(selected_runs, config["key_metrics"])
    image_dpi_percent = resolve_image_dpi_percent(
        config["image_dpi_percent"],
        image_dpi_percent_override,
    )

    output_path = ensure_output_suffix(output_path.expanduser().resolve(), config["output_format"])
    output_path.parent.mkdir(parents=True, exist_ok=True)

    page_number = 1
    writer = make_report_writer(output_path, config["output_format"])
    try:
        page_number = render_image_pages(
            writer=writer,
            selected_runs=selected_runs,
            image_keys=image_keys,
            max_rows=config["max_rows"],
            max_columns=config["max_columns"],
            run_name_max_chars_per_line=config["run_name_max_chars_per_line"],
            image_max_side=image_max_side,
            image_dpi_percent=image_dpi_percent,
            dpi=dpi,
            page_number=page_number,
        )
        page_number = render_metric_pages(
            writer=writer,
            selected_runs=selected_runs,
            metrics=metrics,
            dpi=dpi,
            page_number=page_number,
        )
        render_hyperparameter_page(
            writer=writer,
            selected_runs=selected_runs,
            hyperparameter_names=config["key_hyperparameters"],
            run_name_max_chars_per_line=config["run_name_max_chars_per_line"],
            dpi=dpi,
            page_number=page_number,
        )
    finally:
        writer.close()

    return output_path


def main() -> int:
    args = parse_args()
    try:
        output_path = args.output if args.output is not None else args.output_pdf
        saved_report = build_report(
            config_path=args.config.resolve(),
            output_path=output_path,
            dpi=args.dpi,
            image_max_side=args.image_max_side,
            image_dpi_percent_override=args.image_dpi_percent,
        )
    except ConfigError as exc:
        print(f"Config error: {exc}")
        return 2
    except Exception as exc:
        print(f"Failed to build report: {exc}")
        return 1

    print(f"Saved report to {saved_report}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
